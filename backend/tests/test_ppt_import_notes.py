"""P0 fix (2026-02): PPT Import must populate text into Presenter Notes.

Originally the high-fidelity PPT parser only set `slide.notes` when the
source PPTX had explicit speaker notes. For PPTs that did NOT have
presenter notes the panel was empty even though there was plenty of
text inside the slide shapes, so the author had no context to work with.

This regression test covers four scenarios:
  1. Real PPT presenter notes are preserved verbatim.
  2. PPT without notes: slide body text is copied into `notes`.
  3. PPT with whitespace-only notes: treated as empty -> body text fallback.
  4. PPT without notes and without body text: notes stays empty (no crash).

It also verifies that `extractedText` is always populated as a separate
field so the AI Tutor / SCORM exporter (which read `extractedText`
independently) keep working even if the author later edits notes.
"""
import os
import sys
import shutil
import tempfile
from pathlib import Path
from unittest.mock import patch

import pytest
from PIL import Image
from pptx import Presentation

# Allow importing `services.*` and `models` from the backend root.
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from services import ppt_image_parser as pim  # noqa: E402
from services.ppt_parser import parse_pptx  # noqa: E402


def _make_pptx() -> str:
    prs = Presentation()
    layout = prs.slide_layouts[1]

    # A) Real presenter notes (must win over body text)
    a = prs.slides.add_slide(layout)
    a.shapes.title.text = "Slide A"
    a.placeholders[1].text = "Body A bullet text"
    a.notes_slide.notes_text_frame.text = "REAL PRESENTER NOTES A"

    # B) No presenter notes -> fallback to body text
    b = prs.slides.add_slide(layout)
    b.shapes.title.text = "Slide B"
    b.placeholders[1].text = "Body B line 1\nBody B line 2"

    # C) Whitespace-only notes -> treated as empty, fallback to body text
    c = prs.slides.add_slide(layout)
    c.shapes.title.text = "Slide C"
    c.placeholders[1].text = "Body C content"
    c.notes_slide.notes_text_frame.text = "   "

    # D) No notes, no body text (title only) -> notes stays empty
    d = prs.slides.add_slide(layout)
    d.shapes.title.text = "Slide D"

    f = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir="/tmp")
    prs.save(f.name)
    return f.name


def _fake_image_conversion(output_dir: str, n: int) -> list:
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    paths = []
    for i in range(n):
        p = out / f"slide_{i + 1:03d}.png"
        Image.new("RGB", (1920, 1080), "#FFFFFF").save(str(p))
        paths.append(str(p))
    return paths


@pytest.fixture
def test_pptx_path():
    p = _make_pptx()
    yield p
    try:
        os.unlink(p)
    except OSError:
        pass


def test_high_fidelity_notes_extraction(test_pptx_path, tmp_path):
    """High-fidelity parser must populate notes (real or fallback)."""

    def fake_convert(pptx_path, output_dir, **kwargs):
        return _fake_image_conversion(output_dir, n=4)

    with patch.object(pim, "convert_pptx_to_images", side_effect=fake_convert):
        course = pim.parse_pptx_high_fidelity(test_pptx_path, "pid", str(tmp_path))

    assert len(course.slides) == 4

    s1, s2, s3, s4 = course.slides

    # 1) Real notes preserved
    assert s1.notes == "REAL PRESENTER NOTES A"
    assert "Body A bullet text" in (s1.extractedText or "")

    # 2) No notes -> fallback to body text
    assert s2.notes is not None
    assert "Body B line 1" in s2.notes
    assert "Body B line 2" in s2.notes
    assert s2.extractedText == s2.notes

    # 3) Whitespace notes treated as empty -> fallback
    assert s3.notes is not None
    assert "Body C content" in s3.notes
    assert s3.extractedText == s3.notes

    # 4) No notes, no body -> notes stays empty
    assert not s4.notes
    assert not s4.extractedText


def test_python_only_parser_notes_extraction(test_pptx_path, tmp_path):
    """Fallback (Python-only) parser must also populate notes + extractedText."""
    course = parse_pptx(test_pptx_path, "pid", str(tmp_path))

    assert len(course.slides) == 4
    s1, s2, s3, s4 = course.slides

    assert s1.notes == "REAL PRESENTER NOTES A"
    assert "Body A bullet text" in (s1.extractedText or "")

    assert s2.notes and "Body B line 1" in s2.notes
    assert s2.extractedText and "Body B line 1" in s2.extractedText

    assert s3.notes and "Body C content" in s3.notes

    assert not s4.notes
    assert not s4.extractedText


def test_serialized_payload_contains_both_fields(test_pptx_path, tmp_path):
    """`model_dump()` (what gets persisted to Mongo) must contain both keys."""
    course = parse_pptx(test_pptx_path, "pid", str(tmp_path))
    payload = course.model_dump()

    for i, sl in enumerate(payload["slides"]):
        assert "notes" in sl, f"slide {i + 1} missing 'notes' key in payload"
        assert "extractedText" in sl, f"slide {i + 1} missing 'extractedText' key"
