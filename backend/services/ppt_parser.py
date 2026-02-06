"""
PPT/PPTX Parser Service
Extracts slides, elements, animations, and transitions from PowerPoint files
"""
import os
import io
import base64
import logging
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import uuid

from models import (
    Course, CourseMetadata, Slide, SlideElement, ElementStyle,
    Animation, SlideTransition, SlideAudio
)

logger = logging.getLogger(__name__)

# EMU to pixels conversion (96 DPI)
EMU_PER_INCH = 914400
DPI = 96

def emu_to_px(emu: int) -> float:
    """Convert EMU to pixels"""
    if emu is None:
        return 0
    return (emu / EMU_PER_INCH) * DPI

def rgb_to_hex(rgb) -> str:
    """Convert RGBColor to hex string"""
    if rgb is None:
        return None
    try:
        if isinstance(rgb, RGBColor):
            return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}".upper()
        return None
    except:
        return None

def get_fill_color(shape) -> Optional[str]:
    """Extract fill color from shape"""
    try:
        if hasattr(shape, 'fill'):
            fill = shape.fill
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    try:
                        rgb = fill.fore_color.rgb
                        if rgb:
                            return rgb_to_hex(rgb)
                    except:
                        pass
        return None
    except:
        return None

def get_line_color(shape) -> Optional[str]:
    """Extract line/stroke color from shape"""
    try:
        if hasattr(shape, 'line'):
            line = shape.line
            if line.fill and line.fill.fore_color:
                try:
                    rgb = line.fill.fore_color.rgb
                    if rgb:
                        return rgb_to_hex(rgb)
                except:
                    pass
        return None
    except:
        return None

def extract_text_properties(paragraph) -> Dict[str, Any]:
    """Extract text properties from a paragraph"""
    props = {}
    try:
        # Alignment
        if paragraph.alignment:
            align_map = {
                PP_ALIGN.LEFT: 'left',
                PP_ALIGN.CENTER: 'center',
                PP_ALIGN.RIGHT: 'right',
                PP_ALIGN.JUSTIFY: 'justify'
            }
            props['textAlign'] = align_map.get(paragraph.alignment, 'left')
        
        # Get first run properties
        if paragraph.runs:
            run = paragraph.runs[0]
            font = run.font
            
            if font.name:
                props['fontFamily'] = font.name
            if font.size:
                props['fontSize'] = font.size.pt
            if font.bold:
                props['fontWeight'] = 'bold'
            if font.color and font.color.rgb:
                props['fontColor'] = rgb_to_hex(font.color.rgb)
    except Exception as e:
        logger.debug(f"Error extracting text props: {e}")
    
    return props

def extract_smartart(shape, assets_dir: Path, project_id: str) -> Optional[Tuple[str, str, List[str]]]:
    """
    Extract SmartArt graphic and convert to image
    Returns: (image_url, smartart_type, text_content_list) or None
    """
    try:
        from pptx.oxml.ns import qn
        
        shape_xml = shape._element
        
        # Check if this shape contains a SmartArt graphic (dgm:relIds)
        # SmartArt is typically in graphicFrame elements with dgm namespace
        dgm_rels = shape_xml.findall('.//' + qn('dgm:relIds'))
        
        if not dgm_rels:
            # Also check for graphicData with SmartArt
            graphic_data = shape_xml.find('.//' + qn('a:graphicData'))
            if graphic_data is not None:
                uri = graphic_data.get('uri', '')
                if 'dgm' not in uri.lower() and 'smartart' not in uri.lower():
                    return None
            else:
                return None
        
        logger.info(f"Found SmartArt graphic")
        
        # Extract text content from SmartArt for accessibility
        text_content = []
        
        # SmartArt text is often in t elements
        text_elements = shape_xml.findall('.//' + qn('a:t'))
        for t_elem in text_elements:
            if t_elem.text:
                text_content.append(t_elem.text.strip())
        
        # Also try to find text in dgm:t elements
        dgm_texts = shape_xml.findall('.//' + qn('dgm:t'))
        for dgm_t in dgm_texts:
            if dgm_t.text:
                text_content.append(dgm_t.text.strip())
        
        # Try to determine SmartArt type
        smartart_type = "diagram"
        
        # Look for layout definition
        layout_node = shape_xml.find('.//' + qn('dgm:layoutNode'))
        if layout_node is not None:
            name = layout_node.get('name', '')
            if name:
                smartart_type = name.lower()
        
        # Try to extract embedded image representation
        blips = shape_xml.findall('.//' + qn('a:blip'))
        for blip in blips:
            embed_id = blip.get(qn('r:embed'))
            if embed_id:
                try:
                    part = shape.part
                    image_part = part.related_part(embed_id)
                    if image_part:
                        image_bytes = image_part.blob
                        content_type = image_part.content_type
                        ext_map = {
                            'image/png': 'png',
                            'image/jpeg': 'jpg',
                            'image/gif': 'gif',
                            'image/emf': 'emf',
                            'image/wmf': 'wmf'
                        }
                        ext = ext_map.get(content_type, 'png')
                        image_filename = f"smartart_{uuid.uuid4()}.{ext}"
                        image_path = assets_dir / image_filename
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        
                        # Convert EMF/WMF to PNG if needed
                        if ext in ['emf', 'wmf']:
                            try:
                                from PIL import Image
                                png_filename = f"smartart_{uuid.uuid4()}.png"
                                png_path = assets_dir / png_filename
                                
                                # Try to convert using PIL (may not work for all EMF)
                                with Image.open(image_path) as img:
                                    img.save(png_path, 'PNG')
                                
                                # Use the PNG version
                                os.remove(image_path)
                                image_filename = png_filename
                            except Exception as e:
                                logger.debug(f"Could not convert EMF/WMF to PNG: {e}")
                        
                        image_url = f"/api/projects/{project_id}/assets/{image_filename}"
                        logger.info(f"Extracted SmartArt image: {image_filename}")
                        return (image_url, smartart_type, text_content)
                        
                except Exception as e:
                    logger.debug(f"Error extracting SmartArt image: {e}")
        
        # If no embedded image, return None (will need fallback rendering)
        return (None, smartart_type, text_content)
        
    except Exception as e:
        logger.debug(f"Error extracting SmartArt: {e}")
        return None


def shape_to_element(shape, index: int, assets_dir: Path, project_id: str) -> Optional[SlideElement]:
    """Convert a PowerPoint shape to a SlideElement"""
    try:
        element_id = str(uuid.uuid4())
        
        # Basic positioning
        x = emu_to_px(shape.left) if shape.left else 0
        y = emu_to_px(shape.top) if shape.top else 0
        width = emu_to_px(shape.width) if shape.width else 100
        height = emu_to_px(shape.height) if shape.height else 100
        rotation = shape.rotation if hasattr(shape, 'rotation') else 0
        
        style = ElementStyle()
        element_type = "shape"
        content = None
        src = None
        shape_type = None
        svg_path = None
        fallback_image = None
        hyperlink = None
        
        # Check for hyperlink
        if hasattr(shape, 'click_action') and shape.click_action:
            if shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                hyperlink = shape.click_action.hyperlink.address
        
        # Determine shape type
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element_type = "image"
            # Extract image
            try:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_filename = f"{element_id}.{image_ext}"
                image_path = assets_dir / image_filename
                
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                src = f"/api/projects/{project_id}/assets/{image_filename}"
            except Exception as e:
                logger.warning(f"Failed to extract image: {e}")
                return None
                
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or hasattr(shape, 'text_frame'):
            element_type = "text"
            try:
                if hasattr(shape, 'text_frame'):
                    text_frame = shape.text_frame
                    paragraphs = []
                    
                    for para in text_frame.paragraphs:
                        text = para.text
                        if text:
                            paragraphs.append(text)
                            # Extract style from first paragraph
                            if not style.fontFamily:
                                text_props = extract_text_properties(para)
                                for key, val in text_props.items():
                                    setattr(style, key, val)
                    
                    content = '\n'.join(paragraphs)
                    
                    # Vertical alignment
                    if hasattr(text_frame, 'vertical_anchor'):
                        v_align_map = {
                            MSO_ANCHOR.TOP: 'top',
                            MSO_ANCHOR.MIDDLE: 'middle',
                            MSO_ANCHOR.BOTTOM: 'bottom'
                        }
                        style.verticalAlign = v_align_map.get(text_frame.vertical_anchor, 'top')
            except Exception as e:
                logger.debug(f"Error extracting text: {e}")
                
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            element_type = "shape"
            if hasattr(shape, 'auto_shape_type'):
                shape_type = str(shape.auto_shape_type).split('.')[-1].lower()
            
            # Check if this autoshape contains an image (picture fill)
            placeholder_img = extract_placeholder_image(shape, assets_dir, project_id)
            if placeholder_img:
                element_type = "image"
                src = placeholder_img
            
            # Also check for text in shapes
            if hasattr(shape, 'text_frame') and shape.text_frame:
                try:
                    text = shape.text_frame.text
                    if text:
                        content = text
                except:
                    pass
                    
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # For groups, we'd ideally recurse, but for now render as image fallback
            element_type = "shape"
            shape_type = "group"
            
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            element_type = "table"
            # Extract table data
            try:
                if hasattr(shape, 'table'):
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = []
                        for cell in row.cells:
                            cells.append(cell.text)
                        rows.append(cells)
                    content = str(rows)  # JSON-like representation
            except:
                pass
                
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            element_type = "chart"
            # Charts are complex - render as image fallback
            # Try to extract chart image
            chart_img = extract_placeholder_image(shape, assets_dir, project_id)
            if chart_img:
                src = chart_img
                element_type = "image"  # Render as image
            
        elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            element_type = "video"
        
        # Check for SmartArt (can be in different shape types)
        smartart_result = extract_smartart(shape, assets_dir, project_id)
        if smartart_result:
            smartart_url, smartart_type, smartart_texts = smartart_result
            if smartart_url:
                element_type = "smartart"
                src = smartart_url
                content = ' | '.join(smartart_texts) if smartart_texts else None
                # Store SmartArt type in style
                if not style:
                    style = ElementStyle()
                style.smartartType = smartart_type
                logger.info(f"SmartArt extracted: type={smartart_type}, texts={len(smartart_texts)}")
            elif smartart_texts:
                # No image but has text - create text element with SmartArt content
                element_type = "text"
                content = '\n'.join(smartart_texts)
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Placeholder might contain image
            placeholder_img = extract_placeholder_image(shape, assets_dir, project_id)
            if placeholder_img:
                element_type = "image"
                src = placeholder_img
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                element_type = "text"
                try:
                    content = shape.text_frame.text
                except:
                    pass
        
        else:
            # Unknown shape type - try to extract any embedded image
            placeholder_img = extract_placeholder_image(shape, assets_dir, project_id)
            if placeholder_img:
                element_type = "image"
                src = placeholder_img
            
        # Extract fill color
        fill_color = get_fill_color(shape)
        if fill_color:
            style.fill = fill_color
            
        # Extract stroke/line
        stroke_color = get_line_color(shape)
        if stroke_color:
            style.stroke = stroke_color
        
        # Skip empty elements
        if element_type == "text" and not content:
            return None
            
        return SlideElement(
            id=element_id,
            type=element_type,
            x=x,
            y=y,
            width=width,
            height=height,
            rotation=rotation,
            zIndex=index,
            content=content,
            src=src,
            style=style,
            shapeType=shape_type,
            svgPath=svg_path,
            fallbackImage=fallback_image,
            hyperlink=hyperlink
        )
        
    except Exception as e:
        logger.error(f"Error converting shape: {e}")
        return None

def extract_animations(pptx_slide, elements: List[SlideElement]) -> List[Animation]:
    """Extract animations from slide XML"""
    animations = []
    
    try:
        slide_part = pptx_slide.part
        slide_xml = slide_part._element
        
        # Namespace for PowerPoint
        nsmap = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }
        
        # PowerPoint animations are stored in the timing tree (p:timing)
        timing = slide_xml.find('.//p:timing', nsmap)
        if timing is None:
            logger.debug("No timing element found in slide")
            return animations
        
        logger.info(f"Found timing element, extracting animations...")
        
        # Map of shape IDs to element indices
        shape_to_element = {}
        for idx, shape in enumerate(pptx_slide.shapes):
            if hasattr(shape, 'shape_id'):
                shape_to_element[str(shape.shape_id)] = idx
                logger.debug(f"Shape {shape.shape_id} -> element {idx}")
        
        # Animation effect type mapping (based on preset ID)
        # From: https://docs.microsoft.com/en-us/openspecs/office_standards/ms-pptx/
        preset_effect_map = {
            '1': 'appear', '2': 'fly', '3': 'blinds', '4': 'box', '5': 'checkerboard',
            '6': 'circle', '7': 'crawl', '8': 'diamond', '9': 'dissolve', '10': 'fade',
            '11': 'flash_once', '12': 'flip', '13': 'float', '14': 'fly_in', '15': 'fold',
            '16': 'glide', '17': 'grow_and_turn', '18': 'newsflash', '19': 'peek', '20': 'pinwheel',
            '21': 'plus', '22': 'random_bars', '23': 'random_effects', '24': 'rise_up',
            '25': 'shapes', '26': 'sling', '27': 'spinner', '28': 'split', '29': 'stretch',
            '30': 'strips', '31': 'swish', '32': 'swivel', '33': 'wedge', '34': 'wheel',
            '35': 'wipe', '36': 'zoom', '37': 'bounce', '38': 'credits', '39': 'curve_up',
            '40': 'drop', '41': 'elastic', '42': 'float_down', '43': 'float_up',
            '44': 'basic_swivel', '45': 'basic_zoom', '46': 'boomerang', '47': 'basic_bounce',
            '48': 'compress', '49': 'curve_down', '50': 'curve_left', '51': 'curve_right',
            '52': 'expand', '53': 'magnify', '54': 'pulse', '55': 'spin', '56': 'spiral_in',
            '57': 'spiral_out', '58': 'teeter', '59': 'thread', '60': 'unfold',
        }
        
        # Animation type based on presetClass
        animation_type_map = {
            'entr': 'entrance',
            'exit': 'exit',
            'emph': 'emphasis',
            'path': 'motion',
            'mediacall': 'media'
        }
        
        # Find all par (parallel) and seq (sequential) animation containers
        # The structure is: p:timing/p:tnLst/p:par/p:cTn/p:childTnLst/...
        
        # Find all cTn (common time node) elements with presetClass
        all_ctn = timing.findall('.//p:cTn[@presetClass]', nsmap)
        
        for ctn in all_ctn:
            try:
                preset_class = ctn.get('presetClass', 'entr')
                preset_id = ctn.get('presetID', '10')  # Default to fade
                duration = ctn.get('dur')
                delay_val = ctn.get('delay')
                node_type = ctn.get('nodeType', '')
                
                # Find target element
                tgtEl = ctn.find('.//p:tgtEl', nsmap)
                if tgtEl is None:
                    continue
                
                spTgt = tgtEl.find('.//p:spTgt', nsmap)
                if spTgt is None:
                    continue
                
                shape_id = spTgt.get('spid')
                if not shape_id:
                    continue
                
                # Find element index
                elem_idx = shape_to_element.get(shape_id)
                if elem_idx is None:
                    logger.debug(f"Shape {shape_id} not found in element map")
                    continue
                
                if elem_idx >= len(elements):
                    logger.debug(f"Element index {elem_idx} out of range")
                    continue
                
                element = elements[elem_idx]
                
                # Get effect name
                effect = preset_effect_map.get(preset_id, 'fade')
                
                # Calculate duration (in milliseconds, convert to seconds)
                anim_duration = 0.5
                if duration and duration != 'indefinite':
                    try:
                        anim_duration = int(duration) / 1000.0
                    except:
                        pass
                
                # Calculate delay
                anim_delay = 0.0
                if delay_val:
                    try:
                        anim_delay = int(delay_val) / 1000.0
                    except:
                        pass
                
                # Determine trigger type
                trigger = 'afterPrevious'
                if node_type == 'clickEffect':
                    trigger = 'onClick'
                elif node_type == 'withEffect':
                    trigger = 'withPrevious'
                elif node_type == 'afterEffect':
                    trigger = 'afterPrevious'
                
                # Get animation type
                anim_type = animation_type_map.get(preset_class, 'entrance')
                
                # Create animation object
                animation = Animation(
                    id=str(uuid.uuid4()),
                    type=anim_type,
                    effect=effect,
                    trigger=trigger,
                    duration=anim_duration,
                    delay=anim_delay,
                    easing='ease'
                )
                
                # Add animation to element and make it visible
                if element.animations is None:
                    element.animations = []
                element.animations.append(animation)
                
                # Mark element as visible since it has animations
                # It will appear with the animation effect
                element.visible = True
                
                animations.append(animation)
                logger.info(f"Extracted animation: {anim_type}/{effect} (duration={anim_duration}s, delay={anim_delay}s) for element {elem_idx}")
                
            except Exception as e:
                logger.debug(f"Error processing animation node: {e}")
                continue
        
        if not animations:
            # Try alternative method - look for animEffect elements directly
            anim_effects = timing.findall('.//p:animEffect', nsmap)
            for anim_effect in anim_effects:
                try:
                    filter_val = anim_effect.get('filter', 'fade')
                    # Extract effect name
                    effect = filter_val.split('(')[0] if '(' in filter_val else filter_val
                    
                    # Find parent cTn for timing info
                    parent = anim_effect.getparent()
                    while parent is not None and parent.tag != '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn':
                        parent = parent.getparent()
                    
                    duration = 0.5
                    delay = 0.0
                    if parent is not None:
                        dur = parent.get('dur')
                        if dur and dur != 'indefinite':
                            try:
                                duration = int(dur) / 1000.0
                            except:
                                pass
                    
                    # Try to find target
                    tgtEl = anim_effect.find('.//p:tgtEl', nsmap)
                    if tgtEl is not None:
                        spTgt = tgtEl.find('.//p:spTgt', nsmap)
                        if spTgt is not None:
                            shape_id = spTgt.get('spid')
                            elem_idx = shape_to_element.get(shape_id)
                            if elem_idx is not None and elem_idx < len(elements):
                                element = elements[elem_idx]
                                
                                animation = Animation(
                                    id=str(uuid.uuid4()),
                                    type='entrance',
                                    effect=effect,
                                    trigger='afterPrevious',
                                    duration=duration,
                                    delay=delay,
                                    easing='ease'
                                )
                                
                                if element.animations is None:
                                    element.animations = []
                                element.animations.append(animation)
                                
                                # Mark element as visible since it has animations
                                element.visible = True
                                
                                animations.append(animation)
                                logger.info(f"Extracted animEffect: {effect} for element {elem_idx}")
                except Exception as e:
                    logger.debug(f"Error in animEffect extraction: {e}")
                    continue
        
        logger.info(f"Total animations extracted: {len(animations)}")
        
    except Exception as e:
        logger.error(f"Error extracting animations: {e}")
        import traceback
        traceback.print_exc()
    
    return animations

def extract_transition(pptx_slide) -> SlideTransition:
    """Extract slide transition"""
    transition = SlideTransition()
    
    try:
        # Transitions are in the slide XML under p:transition
        slide_part = pptx_slide.part
        slide_xml = slide_part._element
        
        # Find transition element
        trans_el = slide_xml.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        if trans_el is not None:
            # Get transition type from child elements
            for child in trans_el:
                tag = child.tag.split('}')[-1]
                if tag not in ['sndAc']:  # Skip sound action
                    transition.type = tag.lower()
                    # Get direction if available
                    if 'dir' in child.attrib:
                        transition.direction = child.attrib['dir']
                    break
            
            # Get duration
            if 'spd' in trans_el.attrib:
                speed_map = {'slow': 1.0, 'med': 0.5, 'fast': 0.25}
                transition.duration = speed_map.get(trans_el.attrib['spd'], 0.5)
                
    except Exception as e:
        logger.debug(f"Error extracting transition: {e}")
    
    return transition

def get_slide_background(pptx_slide, assets_dir: Path, project_id: str) -> Tuple[str, Optional[str]]:
    """Extract slide background color or image"""
    bg_color = "#FFFFFF"
    bg_image = None
    
    try:
        # Try to get background from slide
        background = pptx_slide.background
        if background and background.fill:
            fill = background.fill
            
            # Solid fill
            if fill.type is not None:
                try:
                    if hasattr(fill, 'fore_color') and fill.fore_color:
                        rgb = fill.fore_color.rgb
                        if rgb:
                            bg_color = rgb_to_hex(rgb)
                except:
                    pass
            
            # Check for picture fill using XML
            try:
                from pptx.oxml.ns import qn
                bg_xml = background._element
                
                # Look for blipFill in the background
                blip_fills = bg_xml.findall('.//' + qn('a:blip'))
                for blip in blip_fills:
                    embed_id = blip.get(qn('r:embed'))
                    if embed_id:
                        # Get the related part
                        related_part = pptx_slide.part.related_part(embed_id)
                        if related_part:
                            image_bytes = related_part.blob
                            # Determine extension from content type
                            content_type = related_part.content_type
                            ext_map = {
                                'image/png': 'png',
                                'image/jpeg': 'jpg',
                                'image/gif': 'gif',
                                'image/bmp': 'bmp',
                                'image/tiff': 'tiff'
                            }
                            ext = ext_map.get(content_type, 'png')
                            image_filename = f"bg_{uuid.uuid4()}.{ext}"
                            image_path = assets_dir / image_filename
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            bg_image = f"/api/projects/{project_id}/assets/{image_filename}"
                            logger.info(f"Extracted background image: {image_filename}")
                            break
            except Exception as e:
                logger.debug(f"Error extracting blip background: {e}")
                    
    except Exception as e:
        logger.debug(f"Error extracting background: {e}")
    
    return bg_color, bg_image

def extract_placeholder_image(shape, assets_dir: Path, project_id: str) -> Optional[str]:
    """Extract image from placeholder shapes"""
    try:
        from pptx.oxml.ns import qn
        
        # Check if this shape has a blip (embedded image)
        shape_xml = shape._element
        blips = shape_xml.findall('.//' + qn('a:blip'))
        
        for blip in blips:
            embed_id = blip.get(qn('r:embed'))
            if embed_id:
                try:
                    # Get the image part
                    part = shape.part
                    image_part = part.related_part(embed_id)
                    if image_part:
                        image_bytes = image_part.blob
                        content_type = image_part.content_type
                        ext_map = {
                            'image/png': 'png',
                            'image/jpeg': 'jpg', 
                            'image/gif': 'gif',
                            'image/bmp': 'bmp',
                            'image/tiff': 'tiff'
                        }
                        ext = ext_map.get(content_type, 'png')
                        image_filename = f"{uuid.uuid4()}.{ext}"
                        image_path = assets_dir / image_filename
                        
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        
                        logger.info(f"Extracted placeholder image: {image_filename}")
                        return f"/api/projects/{project_id}/assets/{image_filename}"
                except Exception as e:
                    logger.debug(f"Error getting related part: {e}")
    except Exception as e:
        logger.debug(f"Error extracting placeholder image: {e}")
    
    return None

def parse_pptx(file_path: str, project_id: str, storage_dir: str) -> Course:
    """
    Parse a PPTX file and return a Course object
    """
    logger.info(f"Parsing PPTX: {file_path}")
    
    # Create storage directories
    project_dir = Path(storage_dir) / project_id
    assets_dir = project_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    
    conversion_report = {
        "success": [],
        "warnings": [],
        "errors": [],
        "fallbacks": []
    }
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        logger.error(f"Failed to open PPTX: {e}")
        raise ValueError(f"Could not open PowerPoint file: {e}")
    
    # Get presentation dimensions
    slide_width = emu_to_px(prs.slide_width)
    slide_height = emu_to_px(prs.slide_height)
    
    # Extract metadata
    core_props = prs.core_properties
    metadata = CourseMetadata(
        title=core_props.title or Path(file_path).stem,
        author=core_props.author or "",
        description=core_props.subject or "",
        keywords=core_props.keywords.split(',') if core_props.keywords else []
    )
    
    slides = []
    
    for slide_idx, pptx_slide in enumerate(prs.slides):
        slide_id = str(uuid.uuid4())
        
        # Get background
        bg_color, bg_image = get_slide_background(pptx_slide, assets_dir, project_id)
        
        # Get transition
        transition = extract_transition(pptx_slide)
        
        # Extract elements
        elements = []
        for shape_idx, shape in enumerate(pptx_slide.shapes):
            element = shape_to_element(shape, shape_idx, assets_dir, project_id)
            if element:
                elements.append(element)
                conversion_report["success"].append(f"Slide {slide_idx + 1}: {element.type}")
        
        # Extract animations and attach to elements
        slide_animations = extract_animations(pptx_slide, elements)
        if slide_animations:
            conversion_report["success"].append(f"Slide {slide_idx + 1}: {len(slide_animations)} animations")
        
        # Get slide title
        title = f"Slide {slide_idx + 1}"
        if pptx_slide.shapes.title:
            try:
                title = pptx_slide.shapes.title.text or title
            except:
                pass
        
        # Get notes
        notes = None
        if pptx_slide.has_notes_slide:
            try:
                notes_slide = pptx_slide.notes_slide
                notes = notes_slide.notes_text_frame.text
            except:
                pass
        
        slide = Slide(
            id=slide_id,
            title=title[:50],  # Limit title length
            order=slide_idx,
            width=slide_width,
            height=slide_height,
            background=bg_color,
            backgroundImage=bg_image,
            elements=elements,
            transition=transition,
            notes=notes
        )
        
        slides.append(slide)
        logger.info(f"Parsed slide {slide_idx + 1} with {len(elements)} elements")
    
    course = Course(
        metadata=metadata,
        slides=slides,
        originalFilename=Path(file_path).name,
        conversionReport=conversion_report
    )
    
    logger.info(f"Successfully parsed {len(slides)} slides")
    return course

def render_slide_thumbnail(slide: Slide, assets_dir: Path, size: Tuple[int, int] = (320, 180)) -> str:
    """Generate a thumbnail for a slide"""
    try:
        # Create a simple thumbnail using PIL
        img = Image.new('RGB', (int(slide.width), int(slide.height)), slide.background or '#FFFFFF')
        
        # Resize to thumbnail
        img.thumbnail(size, Image.Resampling.LANCZOS)
        
        # Save
        thumb_filename = f"thumb_{slide.id}.png"
        thumb_path = assets_dir / thumb_filename
        img.save(thumb_path, 'PNG')
        
        return f"/assets/{thumb_filename}"
    except Exception as e:
        logger.error(f"Error generating thumbnail: {e}")
        return None
