"""
High-Fidelity PPT Parser using LibreOffice or ConvertAPI
Converts each slide to PNG image for exact visual fidelity
"""
import os
import subprocess
import shutil
import logging
import base64
from pathlib import Path
from typing import List, Tuple, Optional
import uuid
from PIL import Image

from pptx import Presentation
from pptx.util import Emu

from models import (
    Course, CourseMetadata, Slide, SlideElement, ElementStyle,
    Animation, SlideTransition
)
from utils.system_deps import get_libreoffice_path, ensure_system_dependencies
from services.ppt_parser import extract_animations, extract_smartart, extract_animation_masks

logger = logging.getLogger(__name__)

# EMU to pixels conversion
EMU_PER_INCH = 914400
DPI = 96

def emu_to_px(emu: int) -> float:
    if emu is None:
        return 0
    return (emu / EMU_PER_INCH) * DPI


def convert_pptx_to_images_cloud(pptx_path: str, output_dir: str, dpi: int = 150) -> List[str]:
    """
    Convert PPTX to PNG images using ConvertAPI (cloud service).
    This is used when LibreOffice is not available (e.g., in production/Kubernetes).
    Maintains high visual fidelity.
    """
    try:
        import convertapi
    except ImportError:
        logger.error("ConvertAPI not installed. Run: pip install convertapi")
        return []
    
    # Get API key from environment
    api_key = os.environ.get('CONVERTAPI_SECRET')
    if not api_key:
        logger.error("CONVERTAPI_SECRET environment variable not set - cannot use cloud conversion")
        return []
    
    logger.info(f"ConvertAPI key found (length: {len(api_key)})")
    convertapi.api_credentials = api_key
    
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    try:
        logger.info(f"Converting PPTX to PNG using ConvertAPI: {pptx_path}")
        
        # Convert PPTX to PNG with high quality settings
        result = convertapi.convert('png', {
            'File': pptx_path,
            'ImageResolution': str(dpi),
            'ImageQuality': '100',
            'ScaleImage': 'false',
            'ScaleProportions': 'true'
        }, from_format='pptx')
        
        # Get files from result in order (ConvertAPI returns them in slide order)
        image_paths = []
        files = result.files if hasattr(result, 'files') else []
        
        logger.info(f"ConvertAPI returned {len(files)} files")
        
        # Save each file with our naming convention, preserving API order
        for i, file_obj in enumerate(files):
            slide_filename = f"slide_{i+1:03d}.png"
            slide_path = output_path / slide_filename
            
            # Save the file
            file_obj.save(str(slide_path))
            image_paths.append(str(slide_path))
            logger.info(f"Saved slide {i+1}: {slide_filename}")
        
        logger.info(f"ConvertAPI successfully converted {len(image_paths)} slides in order")
        return image_paths
        
    except Exception as e:
        logger.error(f"ConvertAPI conversion failed: {e}")
        return []



def check_slide_has_animations(pptx_slide) -> bool:
    """Check if a slide has any animations defined"""
    try:
        slide_xml = pptx_slide.part._element
        
        # PowerPoint animations are in the timing element
        nsmap = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }
        
        timing = slide_xml.find('.//p:timing', nsmap)
        if timing is None:
            return False
        
        # Check for any cTn (common time node) with presetClass (indicates animation)
        ctn_with_preset = timing.findall('.//p:cTn[@presetClass]', nsmap)
        if ctn_with_preset:
            return True
        
        # Also check for animEffect elements
        anim_effects = timing.findall('.//p:animEffect', nsmap)
        if anim_effects:
            return True
        
        return False
    except Exception as e:
        logger.debug(f"Error checking animations: {e}")
        return False


def convert_pptx_to_images(pptx_path: str, output_dir: str, dpi: int = 150) -> List[str]:
    """
    Convert PPTX to PNG images using LibreOffice (local) or ConvertAPI (cloud).
    Returns list of image paths.
    Tries LibreOffice first, then ConvertAPI as fallback for production environments.
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Create temp directory for conversion
    temp_dir = Path(output_dir) / "temp_convert"
    temp_dir.mkdir(exist_ok=True)
    
    # Find libreoffice executable
    libreoffice_path = get_libreoffice_path()
    
    # If LibreOffice not found, try ConvertAPI (cloud service)
    if not libreoffice_path:
        logger.warning("LibreOffice not found - trying ConvertAPI cloud conversion")
        cloud_result = convert_pptx_to_images_cloud(pptx_path, output_dir, dpi)
        if cloud_result:
            return cloud_result
        logger.warning("ConvertAPI also failed - PPT import will use Python-only parser")
        return []
    
    try:
        # First convert PPTX to PDF using LibreOffice
        # SECURITY (2026-05-15): we used to shell out via `subprocess.run(cmd, shell=True)`
        # with `pptx_path` interpolated into a quoted string. That's vulnerable to
        # command injection if the upload filename contains a `"` (e.g. an attacker
        # uploads `evil"; rm -rf /; #.pptx`). Even though our auth gates filter
        # uploads, defense in depth says NEVER use shell=True with user-supplied
        # paths. Pass arguments as a list so subprocess does the escaping for us.
        cmd_pdf = [
            libreoffice_path,
            "--headless",
            "--invisible",
            "--convert-to", "pdf",
            "--outdir", str(temp_dir),
            str(pptx_path),
        ]
        logger.info(f"Converting PPTX to PDF: {' '.join(cmd_pdf)}")
        result = subprocess.run(cmd_pdf, capture_output=True, text=True, timeout=120)
        
        if result.returncode != 0:
            logger.error(f"LibreOffice PDF conversion failed: {result.stderr}")
            # Try ConvertAPI as fallback
            cloud_result = convert_pptx_to_images_cloud(pptx_path, output_dir, dpi)
            if cloud_result:
                return cloud_result
            return []
        
        # Find the PDF file
        pdf_files = list(temp_dir.glob("*.pdf"))
        if not pdf_files:
            logger.error("No PDF file generated")
            # Try ConvertAPI as fallback
            cloud_result = convert_pptx_to_images_cloud(pptx_path, output_dir, dpi)
            if cloud_result:
                return cloud_result
            return []
        
        pdf_path = pdf_files[0]
        
        # Convert PDF to images using pdftoppm (if available) or PIL
        try:
            # Try pdftoppm for better quality
            from pdf2image import convert_from_path
            images = convert_from_path(str(pdf_path), dpi=dpi)
            
            image_paths = []
            for i, img in enumerate(images):
                img_filename = f"slide_{i+1:03d}.png"
                img_path = output_path / img_filename
                img.save(str(img_path), 'PNG')
                image_paths.append(str(img_path))
                logger.info(f"Saved slide image: {img_filename}")
            
            return image_paths
            
        except ImportError:
            logger.warning("pdf2image not available, using fallback")
            return convert_pptx_to_images_fallback(pptx_path, output_dir)
            
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return convert_pptx_to_images_fallback(pptx_path, output_dir)
    except Exception as e:
        logger.error(f"Error converting PPTX: {e}")
        return convert_pptx_to_images_fallback(pptx_path, output_dir)
    finally:
        # Cleanup temp directory
        if temp_dir.exists():
            shutil.rmtree(temp_dir, ignore_errors=True)

def convert_pptx_to_images_fallback(pptx_path: str, output_dir: str) -> List[str]:
    """
    Fallback: Convert PPTX to images using multiple methods
    1. Try ConvertAPI (cloud service)
    2. Try LibreOffice to PDF, then pdftoppm
    3. Try LibreOffice direct to PNG (only gets first slide)
    Returns empty list if all methods fail (will trigger Python-only parser)
    """
    output_path = Path(output_dir)
    temp_dir = output_path / "temp_fallback"
    temp_dir.mkdir(exist_ok=True)
    
    # First try ConvertAPI (works in production without LibreOffice)
    cloud_result = convert_pptx_to_images_cloud(pptx_path, output_dir)
    if cloud_result:
        return cloud_result
    
    # Find libreoffice executable
    libreoffice_path = get_libreoffice_path()
    
    if not libreoffice_path:
        logger.warning("LibreOffice not found and ConvertAPI failed - returning empty list")
        return []
    
    try:
        # Method 1: Convert to PDF first, then use pdftoppm directly
        # Same security note as the primary path: avoid shell=True with user-supplied paths.
        cmd_pdf = [
            libreoffice_path,
            "--headless",
            "--invisible",
            "--convert-to", "pdf",
            "--outdir", str(temp_dir),
            str(pptx_path),
        ]
        logger.info(f"Fallback: Converting to PDF first: {' '.join(cmd_pdf)}")
        result = subprocess.run(cmd_pdf, capture_output=True, text=True, timeout=120)
        
        pdf_files = list(temp_dir.glob("*.pdf"))
        if pdf_files:
            pdf_path = pdf_files[0]
            
            # Use pdftoppm to convert PDF to images
            output_prefix = str(output_path / "slide")
            cmd_pdftoppm = ["pdftoppm", "-png", "-r", "150", str(pdf_path), output_prefix]
            
            logger.info(f"Converting PDF to images: {' '.join(cmd_pdftoppm)}")
            result = subprocess.run(cmd_pdftoppm, capture_output=True, text=True, timeout=120)
            
            # pdftoppm creates files like slide-1.png, slide-2.png, etc.
            png_files = sorted(output_path.glob("slide-*.png"))
            
            if png_files:
                # Rename to our format: slide_001.png, slide_002.png, etc.
                renamed_files = []
                for i, f in enumerate(png_files):
                    new_name = output_path / f"slide_{i+1:03d}.png"
                    f.rename(new_name)
                    renamed_files.append(str(new_name))
                    logger.info(f"Created slide image: {new_name.name}")
                return renamed_files
        
        # Method 2: Direct PNG (last resort - only first slide)
        logger.warning("PDF method failed, trying direct PNG (will only get first slide)")
        cmd = [
            libreoffice_path,
            "--headless",
            "--invisible",
            "--convert-to", "png",
            "--outdir", str(output_path),
            str(pptx_path),
        ]
        logger.info(f"Fallback conversion: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        png_files = sorted(output_path.glob("*.png"))
        return [str(f) for f in png_files]
        
    except Exception as e:
        logger.error(f"Fallback conversion failed: {e}")
        return []
    finally:
        # Cleanup temp directory
        if temp_dir.exists():
            shutil.rmtree(temp_dir, ignore_errors=True)

def parse_pptx_high_fidelity(file_path: str, project_id: str, storage_dir: str) -> Course:
    """
    Parse PPTX with high visual fidelity by rendering slides as images.
    Falls back to Python-only parsing if LibreOffice is not available.
    """
    logger.info(f"High-fidelity parsing PPTX: {file_path}")
    
    # Create storage directories
    project_dir = Path(storage_dir) / project_id
    assets_dir = project_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    
    conversion_report = {
        "success": [],
        "warnings": [],
        "errors": [],
        "method": "high_fidelity_image"
    }
    
    # Open PPTX for metadata and text extraction
    try:
        prs = Presentation(file_path)
    except Exception as e:
        logger.error(f"Failed to open PPTX: {e}")
        raise ValueError(f"Could not open PowerPoint file: {e}")
    
    # Get presentation dimensions
    slide_width = emu_to_px(prs.slide_width)
    slide_height = emu_to_px(prs.slide_height)
    
    # Extract metadata
    core_props = prs.core_properties
    metadata = CourseMetadata(
        title=core_props.title or Path(file_path).stem,
        author=core_props.author or "",
        description=core_props.subject or "",
        keywords=core_props.keywords.split(',') if core_props.keywords else []
    )
    
    # Try to convert slides to images (LibreOffice local or ConvertAPI cloud)
    logger.info("Attempting to convert slides to images (LibreOffice or ConvertAPI)...")
    slide_images = convert_pptx_to_images(file_path, str(assets_dir))
    
    logger.info(f"convert_pptx_to_images returned {len(slide_images)} images")
    for i, img_path in enumerate(slide_images):
        logger.info(f"  Image {i+1}: {img_path} (exists: {Path(img_path).exists()})")
    
    # Persist slide images in MongoDB for production environments with ephemeral storage
    if slide_images:
        try:
            from services.asset_store import _get_content_type
            mongo_url = os.environ.get('MONGO_URL', '')
            db_name = os.environ.get('DB_NAME', '')
            if mongo_url and db_name:
                from pymongo import MongoClient as _PersistClient
                _is_atlas = "mongodb.net" in mongo_url or "mongodb+srv" in mongo_url
                _pclient = _PersistClient(
                    mongo_url,
                    serverSelectionTimeoutMS=120000 if _is_atlas else 10000,
                    connectTimeoutMS=120000 if _is_atlas else 10000,
                    socketTimeoutMS=600000 if _is_atlas else 30000,
                    retryWrites=True,
                    maxPoolSize=2,
                )
                _pdb = _pclient[db_name]
                persisted = 0
                for img_path in slide_images:
                    img_filename = Path(img_path).name
                    for attempt in range(3):
                        try:
                            with open(img_path, 'rb') as f:
                                data_b64 = base64.b64encode(f.read()).decode('ascii')
                            _pdb.project_assets.update_one(
                                {"project_id": project_id, "filename": img_filename},
                                {"$set": {
                                    "project_id": project_id,
                                    "filename": img_filename,
                                    "data": data_b64,
                                    "content_type": _get_content_type(img_filename),
                                }},
                                upsert=True,
                            )
                            persisted += 1
                            break
                        except Exception as e:
                            logger.warning(f"Failed to persist slide image {img_filename} (attempt {attempt+1}): {e}")
                            if attempt < 2:
                                import time as _time
                                _time.sleep(2)
                    if _is_atlas and persisted > 0 and persisted % 5 == 0:
                        import time as _time
                        _time.sleep(0.5)
                _pclient.close()
                logger.info(f"Persisted {persisted}/{len(slide_images)} slide images in MongoDB")
        except Exception as e:
            logger.warning(f"Failed to persist slide images in MongoDB (non-fatal): {e}")
    
    # If no images generated, fall back to Python-only parser
    if not slide_images:
        logger.warning("Image conversion failed (no LibreOffice or ConvertAPI) - using Python-only parser")
        from services.ppt_parser import parse_pptx
        return parse_pptx(file_path, project_id, storage_dir)
    
    slides = []
    num_slides = len(prs.slides)
    
    for slide_idx, pptx_slide in enumerate(prs.slides):
        slide_id = str(uuid.uuid4())
        elements = []
        
        # First, check if this slide has animations
        has_animations = check_slide_has_animations(pptx_slide)
        logger.info(f"Slide {slide_idx + 1}: has_animations={has_animations}")
        
        # Get slide image path if available - ALWAYS use image background now
        background_image = None
        img_path = None
        img_filename = None
        if slide_idx < len(slide_images):
            img_path = slide_images[slide_idx]
            img_filename = Path(img_path).name
            
            # Resize image to match slide dimensions
            try:
                with Image.open(img_path) as img:
                    target_width = int(slide_width)
                    target_height = int(slide_height)
                    
                    if img.size != (target_width, target_height):
                        img_resized = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
                        img_resized.save(img_path)
                
                background_image = f"/api/projects/{project_id}/assets/{img_filename}"
            except Exception as e:
                logger.warning(f"Error processing slide image: {e}")
        
        if has_animations:
            # NEW APPROACH: Use background image + mask elements for animations
            # The background image shows everything, masks cover animated areas and reveal them with animations
            logger.info(f"Slide {slide_idx + 1}: Using MASK-BASED animations with background image")
            
            # Extract animation masks (positions and timing from XML)
            animation_masks = extract_animation_masks(pptx_slide)
            
            if animation_masks:
                # Get slide background color for masks
                slide_bg_color = "#FFFFFF"
                try:
                    if pptx_slide.background.fill.type is not None:
                        fill = pptx_slide.background.fill
                        if hasattr(fill, 'fore_color') and fill.fore_color and fill.fore_color.rgb:
                            slide_bg_color = f"#{fill.fore_color.rgb}"
                except:
                    pass
                
                # Create animation elements using clip-path approach
                # This works with any background (color or image)
                for mask in animation_masks:
                    is_entrance = mask['animation_type'] == 'entrance'
                    is_exit = mask['animation_type'] == 'exit'
                    is_emphasis = mask['animation_type'] == 'emphasis'
                    
                    if is_entrance:
                        # Entrance animation using clip-path reveal
                        mask_element = SlideElement(
                            id=mask['id'],
                            type="animation_clip",  # New type - uses clip-path
                            x=mask['x'],
                            y=mask['y'],
                            width=mask['width'],
                            height=mask['height'],
                            visible=True,
                            zIndex=1000 + len(elements),
                            style=ElementStyle(
                                opacity=1.0
                            ),
                            animations=[Animation(
                                id=str(uuid.uuid4()),
                                type='entrance',
                                effect=mask['effect'],
                                trigger=mask['trigger'],
                                duration=mask['duration'],
                                delay=mask['start_time'],
                                easing='ease'
                            )]
                        )
                        elements.append(mask_element)
                        
                    elif is_exit:
                        # Exit animation using clip-path hide
                        mask_element = SlideElement(
                            id=mask['id'],
                            type="animation_clip",
                            x=mask['x'],
                            y=mask['y'],
                            width=mask['width'],
                            height=mask['height'],
                            visible=True,
                            zIndex=1000 + len(elements),
                            style=ElementStyle(
                                opacity=1.0
                            ),
                            animations=[Animation(
                                id=str(uuid.uuid4()),
                                type='exit',
                                effect=mask['effect'],
                                trigger=mask['trigger'],
                                duration=mask['duration'],
                                delay=mask['start_time'],
                                easing='ease'
                            )]
                        )
                        elements.append(mask_element)
                        
                    elif is_emphasis:
                        # Emphasis: create a highlight overlay effect
                        mask_element = SlideElement(
                            id=mask['id'],
                            type="animation_highlight",
                            x=mask['x'],
                            y=mask['y'],
                            width=mask['width'],
                            height=mask['height'],
                            visible=True,
                            zIndex=1000 + len(elements),
                            style=ElementStyle(
                                fill="transparent",
                                opacity=0.0
                            ),
                            animations=[Animation(
                                id=str(uuid.uuid4()),
                                type='emphasis',
                                effect=mask['effect'],
                                trigger=mask['trigger'],
                                duration=mask['duration'],
                                delay=mask['start_time'],
                                easing='ease'
                            )]
                        )
                        elements.append(mask_element)
                
                conversion_report["success"].append(f"Slide {slide_idx + 1}: {len(animation_masks)} animation masks with background image")
                logger.info(f"Created {len(animation_masks)} animation masks for slide {slide_idx + 1}")
            else:
                conversion_report["success"].append(f"Slide {slide_idx + 1}: Rendered as image (animations detected but no masks extracted)")
        else:
            # For slides WITHOUT animations: just use the image background
            conversion_report["success"].append(f"Slide {slide_idx + 1}: Rendered as image (no animations)")
            logger.info(f"Slide {slide_idx + 1}: No animations, using image only")
        
        # Get slide title
        title = f"Slide {slide_idx + 1}"
        if pptx_slide.shapes.title:
            try:
                title = pptx_slide.shapes.title.text or title
            except:
                pass
        
        # Extract ALL text from slide shapes (critical for narration when images are unavailable)
        slide_text_parts = []
        try:
            for shape in pptx_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text and para_text != title:
                            slide_text_parts.append(para_text)
                # Also extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                slide_text_parts.append(cell_text)
        except Exception as txt_err:
            logger.warning(f"Failed to extract text from slide {slide_idx + 1}: {txt_err}")
        
        extracted_text = "\n".join(slide_text_parts) if slide_text_parts else None
        
        # Get notes
        notes = None
        if pptx_slide.has_notes_slide:
            try:
                notes_slide = pptx_slide.notes_slide
                notes = notes_slide.notes_text_frame.text
            except:
                pass
        
        slide = Slide(
            id=slide_id,
            title=title[:50],
            order=slide_idx,
            width=slide_width,
            height=slide_height,
            background="#FFFFFF",
            backgroundImage=background_image,
            elements=elements,
            notes=notes,
            extractedText=extracted_text,
        )
        
        slides.append(slide)
        logger.info(f"Processed slide {slide_idx + 1}")
    
    course = Course(
        metadata=metadata,
        slides=slides,
        originalFilename=Path(file_path).name,
        conversionReport=conversion_report
    )
    
    logger.info(f"Successfully processed {len(slides)} slides with high fidelity")
    return course
